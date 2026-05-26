"""Build optimization/mean_variance_optimization.pptx.

Theory + this repo's ZAR currency frontier.
"""
from pathlib import Path

from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _build_timing_xml(shape_groups):
    """One <p:par> click block per group; each group can hold several shapes
    that all appear together on the same click."""
    cur_id = 3  # ids 1 and 2 are reserved for tmRoot / mainSeq
    click_blocks = []
    build_entries = []
    seen_spids = set()
    for group in shape_groups:
        id1, id2 = cur_id, cur_id + 1
        cur_id += 2
        effect_blocks = []
        for i, spid in enumerate(group):
            idA, idB = cur_id, cur_id + 1
            cur_id += 2
            node_type = "clickEffect" if i == 0 else "withEffect"
            effect_blocks.append(
                f'<p:par>'
                f'<p:cTn id="{idA}" presetID="1" presetClass="entr" '
                f'presetSubtype="0" fill="hold" grpId="0" nodeType="{node_type}">'
                f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
                f'<p:childTnLst>'
                f'<p:set>'
                f'<p:cBhvr>'
                f'<p:cTn id="{idB}" dur="1" fill="hold">'
                f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
                f'</p:cTn>'
                f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
                f'<p:attrNameLst><p:attrName>style.visibility</p:attrName>'
                f'</p:attrNameLst>'
                f'</p:cBhvr>'
                f'<p:to><p:strVal val="visible"/></p:to>'
                f'</p:set>'
                f'</p:childTnLst>'
                f'</p:cTn>'
                f'</p:par>'
            )
            if spid not in seen_spids:
                build_entries.append(f'<p:bldP spid="{spid}" grpId="0"/>')
                seen_spids.add(spid)
        click_blocks.append(
            f'<p:par>'
            f'<p:cTn id="{id1}" fill="hold">'
            f'<p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>'
            f'<p:childTnLst>'
            f'<p:par>'
            f'<p:cTn id="{id2}" fill="hold">'
            f'<p:stCondLst><p:cond delay="0"/></p:stCondLst>'
            f'<p:childTnLst>{"".join(effect_blocks)}</p:childTnLst>'
            f'</p:cTn>'
            f'</p:par>'
            f'</p:childTnLst>'
            f'</p:cTn>'
            f'</p:par>'
        )
    return (
        f'<p:timing xmlns:p="{P_NS}">'
        f'<p:tnLst>'
        f'<p:par>'
        f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
        f'<p:childTnLst>'
        f'<p:seq concurrent="1" nextAc="seek">'
        f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
        f'<p:childTnLst>{"".join(click_blocks)}</p:childTnLst>'
        f'</p:cTn>'
        f'<p:prevCondLst><p:cond evt="onPrev" delay="0">'
        f'<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
        f'<p:nextCondLst><p:cond evt="onNext" delay="0">'
        f'<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
        f'</p:seq>'
        f'</p:childTnLst>'
        f'</p:cTn>'
        f'</p:par>'
        f'</p:tnLst>'
        f'<p:bldLst>{"".join(build_entries)}</p:bldLst>'
        f'</p:timing>'
    )


def attach_click_animation(slide, shape_groups):
    """shape_groups is a list of lists of Shape objects — each inner list
    appears together on one mouse click, in order."""
    spid_groups = [[str(sh.shape_id) for sh in g] for g in shape_groups]
    xml = _build_timing_xml(spid_groups)
    timing_el = etree.fromstring(xml)
    sld = slide._element
    for old in sld.findall(qn('p:timing')):
        sld.remove(old)
    transition = sld.find(qn('p:transition'))
    clr_map = sld.find(qn('p:clrMapOvr'))
    if transition is not None:
        transition.addnext(timing_el)
    elif clr_map is not None:
        clr_map.addnext(timing_el)
    else:
        sld.append(timing_el)


def add_bullets_separate(slide, left_in, top_in, width_in, items, *,
                         size=16, color=None, gap_in=0.12):
    """Each bullet in its own textbox so it can be animated individually.
    Returns the list of textbox shapes in order."""
    if color is None:
        color = NAVY
    shapes = []
    y = top_in
    chars_per_line = max(1, int(width_in * 11 * (16 / size)))
    for item in items:
        lines = max(1, -(-len(item) // chars_per_line))
        h = 0.32 * lines + 0.10
        tb = slide.shapes.add_textbox(Inches(left_in), Inches(y),
                                      Inches(width_in), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_top = Pt(2)
        tf.margin_bottom = Pt(2)
        p = tf.paragraphs[0]
        p.text = "•  " + item
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = "Calibri"
        shapes.append(tb)
        y += h + gap_in
    return shapes

HERE = Path(__file__).resolve().parent
OUT = HERE / "mean_variance_optimization.pptx"
FRONTIER_IMG = HERE / "zar_debt_frontier.png"
FRONTIER_DOM_IMG = HERE / "zar_debt_frontier_with_domestic.png"

NAVY = RGBColor(0x0B, 0x2E, 0x4F)
ACCENT = RGBColor(0xC8, 0x4B, 0x31)
GREY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF2, 0xF3, 0xF5)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

SW = prs.slide_width
SH = prs.slide_height


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return s


def add_text(slide, left, top, width, height, text, *, size=18,
             bold=False, color=NAVY, align=None):
    from pptx.enum.text import PP_ALIGN
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def add_bullets(slide, left, top, width, height, items, *, size=20,
                color=NAVY):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + item
        p.space_after = Pt(8)
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = "Calibri"
    return tb


def add_title_bar(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.0))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.55),
             title, size=28, bold=True,
             color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.62), Inches(12), Inches(0.35),
                 subtitle, size=14,
                 color=RGBColor(0xCC, 0xD6, 0xE0))


# --- Slide 1: Title -----------------------------------------------------
s = add_slide()
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.6), SW, Inches(2.3))
band.line.fill.background()
band.fill.solid()
band.fill.fore_color.rgb = NAVY
add_text(s, Inches(0.6), Inches(2.85), Inches(12), Inches(0.9),
         "Mean Variance Optimization", size=48, bold=True,
         color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.6), Inches(3.8), Inches(12), Inches(0.6),
         "Theory, and an application to sovereign debt currency choice",
         size=22, color=RGBColor(0xCC, 0xD6, 0xE0))
add_text(s, Inches(0.6), Inches(5.6), Inches(12), Inches(0.4),
         "ZAR debt frontier  —  wb-debt-simulation / optimization",
         size=14, color=GREY)

# --- Slide 2: What is MVO ----------------------------------------------
s = add_slide()
add_title_bar(s, "What is Mean Variance Optimization?",
              "Markowitz (1952) — portfolio selection as a quadratic program")
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5), [
    "Pick portfolio weights x that trade off expected return against risk.",
    "Risk is measured by variance of portfolio return:  xᵀ Σ x.",
    "Expected return is linear in weights:  μᵀ x.",
    "Solutions sweep a risk-aversion parameter λ to trace the efficient frontier.",
    "The same machinery applies to debt-cost minimization — just flip the sign of the linear term.",
])

# --- Slide 3: Formulation ----------------------------------------------
s = add_slide()
add_title_bar(s, "The optimization problem",
              "Quadratic objective, linear constraints — solved as a QP")
add_text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.5),
         "Asset allocation form (return maximization):", size=20, bold=True)
add_text(s, Inches(1.0), Inches(2.0), Inches(11), Inches(0.6),
         "minimize    (1 − λ) · xᵀ Σ x  −  λ · μᵀ x",
         size=22, color=ACCENT, bold=True)
add_text(s, Inches(0.6), Inches(2.9), Inches(12), Inches(0.5),
         "Subject to:", size=20, bold=True)
add_bullets(s, Inches(1.0), Inches(3.4), Inches(11), Inches(2.5), [
    "Σ xᵢ  =  1   (fully invested)",
    "bounds_min  ≤  x  ≤  bounds_max   (box constraints, e.g. no short, sector caps)",
    "Optional linear inequalities  wᵀ x  ≤  c   (e.g. domestic-share floor)",
], size=18)
add_text(s, Inches(0.6), Inches(5.9), Inches(12), Inches(0.5),
         "Cost-minimization form (debt frontier): swap the sign of the linear term —  + λ · μᵀ x.",
         size=16, color=GREY)

# --- Slide 4: Efficient frontier (liability framing) -------------------
s = add_slide()
slide4 = s
add_title_bar(s, "The efficient frontier — debt-cost framing",
              "Sweep λ ∈ [0, 1] → one optimal funding mix per risk-aversion level")
slide4_bullets = add_bullets_separate(s, 0.6, 1.4, 7, [
    "λ = 0:  pure risk minimization  —  the minimum-variance funding mix.",
    "λ = 1:  pure cost minimization  —  cheapest mix, ignores FX risk.",
    "Intermediate λ:  Pareto-optimal trade-offs between debt cost and FX volatility.",
    "Plot (risk, expected cost) for each solution to trace the lower envelope of feasible mixes.",
    "Two-fund property:  every point on the frontier is a linear combination of any two other frontier points — pick A and B, sweep α ∈ [0, 1], and α·A + (1−α)·B traces the curve between them.",
    "Current funding mix sits above the frontier — the gap shows the cost or risk reduction available.",
], size=16)
# small frontier sketch on the right — DEBT version: descending convex curve
left, top, w, h = Inches(8.0), Inches(1.6), Inches(4.8), Inches(5.0)
box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
box.line.color.rgb = GREY
box.fill.solid()
box.fill.fore_color.rgb = LIGHT
# axes
ax_left = left + Inches(0.5)
ax_bot = top + h - Inches(0.6)
ax_right = left + w - Inches(0.3)
ax_top = top + Inches(0.4)
s.shapes.add_connector(1, ax_left, ax_bot, ax_right, ax_bot).line.color.rgb = NAVY
s.shapes.add_connector(1, ax_left, ax_bot, ax_left, ax_top).line.color.rgb = NAVY
# frontier: descending convex curve from upper-left (low risk, high cost)
# to lower-right (high risk, low cost). y_frac = (1 - t)^2.
pts = []
for i in range(40):
    t = i / 39
    x_frac = 0.05 + 0.9 * t            # leave a small left margin
    y_frac = (1 - t) ** 2
    px = ax_left + Emu(int((ax_right - ax_left) * x_frac))
    py = ax_bot - Emu(int((ax_bot - ax_top) * (0.15 + 0.75 * y_frac)))
    pts.append((px, py))
for i in range(len(pts) - 1):
    c = s.shapes.add_connector(1, pts[i][0], pts[i][1],
                               pts[i + 1][0], pts[i + 1][1])
    c.line.color.rgb = ACCENT
    c.line.width = Pt(2.5)
# annotate the two end-points and the current-mix marker
mv_x, mv_y = pts[0]
mc_x, mc_y = pts[-1]
add_text(s, mv_x + Inches(0.05), mv_y - Inches(0.35), Inches(1.6), Inches(0.3),
         "λ = 0  min-variance", size=10, color=NAVY, bold=True)
add_text(s, mc_x - Inches(1.7), mc_y - Inches(0.05), Inches(1.7), Inches(0.3),
         "λ = 1  min-cost", size=10, color=NAVY, bold=True, align="right")

# Two-fund points: A near upper-left, B near lower-right of the frontier
def _marker(slide, cx, cy, label, label_dx, label_dy):
    r = Inches(0.10)
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               cx - r, cy - r, r * 2, r * 2)
    d.line.color.rgb = ACCENT
    d.fill.solid()
    d.fill.fore_color.rgb = ACCENT
    add_text(slide, cx + label_dx, cy + label_dy, Inches(0.6), Inches(0.3),
             label, size=12, color=ACCENT, bold=True)

a_x, a_y = pts[10]      # ~upper-left frontier point
b_x, b_y = pts[29]      # ~lower-right frontier point
_marker(s, a_x, a_y, "A", Inches(-0.32), Inches(-0.30))
_marker(s, b_x, b_y, "B", Inches(0.10),  Inches(0.05))
# current-mix dot, above the frontier
cur_x = ax_left + Emu(int((ax_right - ax_left) * 0.55))
cur_y = ax_bot - Emu(int((ax_bot - ax_top) * 0.55))
dot_r = Inches(0.12)
dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                         cur_x - dot_r, cur_y - dot_r,
                         dot_r * 2, dot_r * 2)
dot.line.color.rgb = NAVY
dot.fill.solid()
dot.fill.fore_color.rgb = NAVY
add_text(s, cur_x + Inches(0.15), cur_y - Inches(0.15), Inches(1.6), Inches(0.3),
         "current mix", size=10, color=NAVY, bold=True)
# axis labels
add_text(s, ax_left + Inches(0.9), ax_bot + Inches(0.05), Inches(2.5), Inches(0.3),
         "Risk  (xᵀ Σ x, FX variance)", size=11, color=GREY)
add_text(s, ax_left - Inches(0.45), ax_top - Inches(0.05), Inches(2.5), Inches(0.3),
         "Expected debt cost", size=11, color=GREY)
add_text(s, ax_left + Inches(1.4), ax_top + Inches(1.1), Inches(2.5), Inches(0.3),
         "efficient frontier", size=12, color=ACCENT, bold=True)

# --- Slide 5: From assets to debt --------------------------------------
s = add_slide()
add_title_bar(s, "From asset allocation to debt management",
              "Same QP, different sign — minimize cost instead of maximizing return")
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5), [
    "A sovereign issuing in multiple currencies faces the same trade-off:",
    "    •  Expected cost  =  interest rate  +  expected FX appreciation against the base currency.",
    "    •  Risk           =  variance of debt-servicing cost driven by FX moves against the base.",
    "Weights x are now funding shares per currency, not asset weights.",
    "Covariance Σ is built from historical log-returns of FX rates against the base (here ZAR).",
    "Set maximize=False so λ tilts the objective toward cost reduction rather than return.",
], size=18)

# --- Slide 6: Pipeline -------------------------------------------------
s = add_slide()
slide6 = s
add_title_bar(s, "The ZAR debt frontier pipeline",
              "optimization/exchangerates_get.py end-to-end")
steps = [
    ("1.  Fetch", "ECB SDMX  →  daily EUR-cross FX  →  cached locally in optimization/currency/"),
    ("2.  Rebase", "convert_base_currency(fx, base='zar')  →  columns become ZAR_USD, ZAR_GBP, ..."),
    ("3.  Returns", "get_fx_returns(...)  —  log-diffs, drop the first NaN row."),
    ("4.  Covariance", "get_fx_covariance(returns)  —  pandas .cov() with a missing-data report."),
    ("5.  Assumptions", "DataFrame: interest_rate, expected_appreciation, min/max_share, current_share."),
    ("6.  Solve", "mv_from_dataframes(cov_df, assumptions, n_points=101)  —  101 QPs along λ."),
    ("7.  Plot", "plot_debt_frontier_labeled(...)  —  cost–risk frontier + funding-share panels."),
]
slide6_step_groups = []
top = Inches(1.3)
for label, body in steps:
    lbl_tb = add_text(s, Inches(0.6), top, Inches(2.2), Inches(0.4),
                      label, size=18, bold=True, color=ACCENT)
    body_tb = add_text(s, Inches(2.9), top, Inches(10), Inches(0.4),
                       body, size=15, color=NAVY)
    slide6_step_groups.append([lbl_tb, body_tb])
    top += Inches(0.65)

# --- Slide 7: Frontier image -------------------------------------------
s = add_slide()
add_title_bar(s, "ZAR debt frontier — worked example",
              "USD, GBP, JPY, CHF, EUR funding for a ZAR-based debt manager")
if FRONTIER_IMG.exists():
    s.shapes.add_picture(str(FRONTIER_IMG), Inches(0.6), Inches(1.25),
                         height=Inches(5.9))
    add_text(s, Inches(0.6), Inches(7.15), Inches(12), Inches(0.3),
             "Source: optimization/zar_debt_frontier.png  —  regenerated by  python exchangerates_get.py",
             size=11, color=GREY)
else:
    add_text(s, Inches(0.6), Inches(3.5), Inches(12), Inches(0.5),
             "[zar_debt_frontier.png not found — run python exchangerates_get.py]",
             size=18, color=ACCENT, align="center")

# --- Slide 8: Adding domestic debt -------------------------------------
s = add_slide()
add_title_bar(s, "Adding domestic debt to the choice set",
              "A ZAR-denominated option:  zero FX risk, 5% interest rate (placeholder)")
if FRONTIER_DOM_IMG.exists():
    s.shapes.add_picture(str(FRONTIER_DOM_IMG), Inches(0.4), Inches(1.2),
                         height=Inches(5.0))
else:
    add_text(s, Inches(0.4), Inches(3.5), Inches(7.5), Inches(0.5),
             "[zar_debt_frontier_with_domestic.png not found]",
             size=14, color=ACCENT, align="center")
# right-hand explainer
add_text(s, Inches(8.6), Inches(1.3), Inches(4.5), Inches(0.5),
         "What changes", size=20, bold=True, color=NAVY)
add_bullets(s, Inches(8.6), Inches(1.8), Inches(4.5), Inches(5.0), [
    "Extend cov_df with a ZAR row/column of zeros — domestic debt has no FX variance.",
    "Add a ZAR row to assumptions:  interest_rate = 5%,  current_share = 0.",
    "The QP now picks among 6 instruments — 5 foreign + 1 domestic.",
    "Min-variance corner shifts to ~100% ZAR (riskless, but expensive).",
    "Min-cost corner is still the cheapest foreign currency, with FX risk.",
    "The widget auto-renders the extra ZAR row — no UI changes needed.",
], size=13)
add_text(s, Inches(0.4), Inches(6.4), Inches(12), Inches(0.3),
         "Source: optimization/zar_debt_frontier_with_domestic.png  —  generated from the notebook’s additional-frontier cell",
         size=11, color=GREY)

# --- Slide 9: Reading the frontier -------------------------------------
s = add_slide()
add_title_bar(s, "Reading the three panels",
              "What the plotter shows you")
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5), [
    "Panel 1  —  cost vs risk:  the efficient frontier, plus the current portfolio marker. If the marker sits below-and-right of the frontier, there is a same-cost / lower-risk move available.",
    "Panel 2  —  funding shares as lines:  each currency’s optimal share as λ moves left (risk-averse) to right (cost-min).",
    "Panel 3  —  stacked area:  the same shares as a 100% composition view — easier to read the mix at any given risk-aversion level.",
    "Row 0 of the result frame is the current composition (off the frontier). Rows 1..n_points are the frontier sweep.",
], size=17)

# --- Slide 9: Inputs in this repo --------------------------------------
s = add_slide()
add_title_bar(s, "DebtFrontierInputs — the editable entry point",
              "Dataclass wrapping covariance + assumptions, with a widget")
add_text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.5),
         "Construct, edit, solve, plot:", size=18, bold=True)
code = [
    "inp = er.DebtFrontierInputs(cov_df=cov_df, assumptions=assumptions,",
    "                            name='basis', chartfolder='graph/')",
    "inp.widget()                  # interactive grid in a notebook",
    "inp.solve()                   # 102-row DataFrame (row 0 = current)",
    "inp.plot()                    # writes graph/basis.svg",
]
tb = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.5), Inches(2.0))
tf = tb.text_frame
tf.word_wrap = True
for i, line in enumerate(code):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    for r in p.runs:
        r.font.size = Pt(16)
        r.font.name = "Consolas"
        r.font.color.rgb = NAVY
add_bullets(s, Inches(0.6), Inches(4.2), Inches(12), Inches(3), [
    "Loop scenarios by mutating inp.assumptions['interest_rate'] and setting inp.name before each plot().",
    "Live Σ current_share indicator in the widget turns green at exactly 1.0.",
    "Default export is SVG only; pass export_formats=('png','pdf','svg') to override.",
], size=16)

# --- Slide 10: Caveats / extensions ------------------------------------
s = add_slide()
add_title_bar(s, "Caveats and extensions", "What to know before pushing this further")
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    "Historical covariance is backward-looking. FX regimes shift — stress with alternative Σ (e.g. crisis-period covariance).",
    "Expected appreciation defaults to 0 in the worked example. Provide a view (forwards, UIP, model output) for production use.",
    "Add constraints via the weights / weigthtedsum slots of mv_opt (note the upstream typo). E.g. domestic-share floor, regional caps.",
    "n_points=101 is fast (<1s for 5 assets). Bump for smoother lines; cost is linear.",
    "Yields here are debt cost, not asset return. Don’t rename the result column — plotter labels are hardcoded.",
    "External dep: mv_opt comes from model_cvx, bundled with ModelFlowIb (not on PyPI).",
], size=16)

# --- Slide 11: Summary -------------------------------------------------
s = add_slide()
add_title_bar(s, "Summary", None)
add_bullets(s, Inches(0.6), Inches(1.5), Inches(12), Inches(5), [
    "MVO is a one-paragraph idea: minimize a quadratic risk term, with a linear return/cost term, under linear constraints.",
    "Sweeping the risk-aversion weight λ traces the efficient frontier.",
    "The debt-manager version is the same QP with the sign flipped on the linear term — cost replaces return.",
    "This repo wires ECB FX data → covariance → QP → labeled frontier plot, with a notebook widget for editing assumptions.",
    "Output: a single chart that makes the cost–risk trade-off across currency mixes legible.",
], size=18)

# --- Animations: slide 4 bullets + slide 6 step rows, on mouse click ----
attach_click_animation(slide4, [[sh] for sh in slide4_bullets])
attach_click_animation(slide6, slide6_step_groups)

prs.save(str(OUT))
print(f"wrote {OUT}")
